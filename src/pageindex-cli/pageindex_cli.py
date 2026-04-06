#!/usr/bin/env python3
"""
PageIndex CLI — Convert PDF, Markdown, DOCX, PPTX, XLSX to JSON node tree
Uses Qwen (or any OpenAI-compatible) API via LiteLLM.
"""

import argparse
import asyncio
import json
import os
import sys
import tempfile
from pathlib import Path

# ── Qwen / LiteLLM env setup (must run before importing litellm/pageindex) ──

def setup_qwen_env(api_key: str | None, api_base: str | None):
    """Configure environment variables for Qwen via LiteLLM OpenAI-compat shim."""
    if api_key:
        os.environ["OPENAI_API_KEY"] = api_key
    if api_base:
        os.environ["OPENAI_API_BASE"] = api_base

# ─────────────────────────────────────────────────────────────────────────────

import litellm
litellm.drop_params = True


# ── Document converters ──────────────────────────────────────────────────────

def convert_docx_to_md(path: str) -> str:
    """Convert a .docx file to Markdown text preserving headings and paragraphs."""
    try:
        from docx import Document
    except ImportError:
        sys.exit("python-docx is required: pip install python-docx")

    doc = Document(path)
    lines = []

    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            lines.append("")
            continue
        if style.startswith("Heading 1"):
            lines.append(f"# {text}")
        elif style.startswith("Heading 2"):
            lines.append(f"## {text}")
        elif style.startswith("Heading 3"):
            lines.append(f"### {text}")
        elif style.startswith("Heading 4"):
            lines.append(f"#### {text}")
        elif style.startswith("Heading 5"):
            lines.append(f"##### {text}")
        elif style.startswith("Heading"):
            lines.append(f"###### {text}")
        else:
            lines.append(text)

    # Also include tables as markdown
    for table in doc.tables:
        if not table.rows:
            continue
        header = [cell.text.strip() for cell in table.rows[0].cells]
        lines.append("")
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in table.rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")

    return "\n".join(lines)


def convert_pptx_to_md(path: str) -> str:
    """Convert a .pptx file to Markdown, treating each slide as a section."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        sys.exit("python-pptx is required: pip install python-pptx")

    prs = Presentation(path)
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        slide_title = ""
        body_parts = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Determine if this shape is the slide title placeholder
            is_title = False
            try:
                pf = shape.placeholder_format
                if pf is not None and pf.idx == 0:
                    is_title = True
            except ValueError:
                pass  # Not a placeholder shape
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if is_title:
                slide_title = text
            else:
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        body_parts.append(para_text)

        heading = slide_title or f"Slide {i}"
        lines.append(f"## {heading}")
        for part in body_parts:
            lines.append(part)
        lines.append("")

    return "\n".join(lines)


def convert_xlsx_to_md(path: str) -> str:
    """Convert a .xlsx file to Markdown with each sheet as a section."""
    try:
        import openpyxl
    except ImportError:
        sys.exit("openpyxl is required: pip install openpyxl")

    wb = openpyxl.load_workbook(path, data_only=True)
    lines = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## {sheet_name}")
        lines.append("")

        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            lines.append("*(empty sheet)*")
            lines.append("")
            continue

        # Find max cols with data
        max_col = max(
            (sum(1 for c in row if c is not None) for row in rows),
            default=0
        )
        if max_col == 0:
            lines.append("*(empty sheet)*")
            lines.append("")
            continue

        # Header row
        header = [str(c) if c is not None else "" for c in rows[0][:max_col]]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * max_col) + " |")

        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row[:max_col]]
            # Pad if shorter
            while len(cells) < max_col:
                cells.append("")
            lines.append("| " + " | ".join(cells) + " |")

        lines.append("")

    return "\n".join(lines)


# ── PageIndex integration ────────────────────────────────────────────────────

def run_pageindex_pdf(pdf_path: str, model: str, opts: dict) -> dict:
    """Run PageIndex on a PDF file and return the tree dict."""
    # Add PageIndex to path
    pageindex_dir = Path(__file__).parent / "PageIndex"
    if str(pageindex_dir) not in sys.path:
        sys.path.insert(0, str(pageindex_dir))

    from pageindex import page_index_main
    from pageindex.utils import ConfigLoader

    user_opt = {k: v for k, v in {
        "model": model,
        "toc_check_page_num": opts.get("toc_check_pages"),
        "max_page_num_each_node": opts.get("max_pages_per_node"),
        "max_token_num_each_node": opts.get("max_tokens_per_node"),
        "if_add_node_id": opts.get("if_add_node_id", "yes"),
        "if_add_node_summary": opts.get("if_add_node_summary", "yes"),
        "if_add_doc_description": opts.get("if_add_doc_description", "yes"),
        "if_add_node_text": opts.get("if_add_node_text", "no"),
    }.items() if v is not None}

    cfg = ConfigLoader().load(user_opt)
    return page_index_main(pdf_path, cfg)


def run_pageindex_md(md_path: str, model: str, opts: dict) -> dict:
    """Run PageIndex on a Markdown file and return the tree dict."""
    pageindex_dir = Path(__file__).parent / "PageIndex"
    if str(pageindex_dir) not in sys.path:
        sys.path.insert(0, str(pageindex_dir))

    from pageindex.page_index_md import md_to_tree

    async def _run():
        return await md_to_tree(
            md_path=md_path,
            if_thinning=opts.get("if_thinning", False),
            min_token_threshold=opts.get("thinning_threshold", 5000),
            if_add_node_summary=opts.get("if_add_node_summary", "yes"),
            summary_token_threshold=opts.get("summary_token_threshold", 200),
            model=model,
            if_add_doc_description=opts.get("if_add_doc_description", "yes"),
            if_add_node_text=opts.get("if_add_node_text", "no"),
            if_add_node_id=opts.get("if_add_node_id", "yes"),
        )

    return asyncio.run(_run())


# ── Main CLI ─────────────────────────────────────────────────────────────────

SUPPORTED_EXTS = {".pdf", ".md", ".markdown", ".docx", ".pptx", ".xlsx"}

BANNER = """
╔═══════════════════════════════════════════════════╗
║          PageIndex CLI  —  Document → JSON        ║
║   PDF · MD · DOCX · PPTX · XLSX → Tree Nodes     ║
╚═══════════════════════════════════════════════════╝
"""

def print_banner():
    print(BANNER)


def detect_format(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext not in SUPPORTED_EXTS:
        sys.exit(
            f"Unsupported file extension '{ext}'.\n"
            f"Supported: {', '.join(sorted(SUPPORTED_EXTS))}"
        )
    return ext


def convert_to_temp_md(input_path: str, ext: str) -> str:
    """Convert office formats to a temporary Markdown file. Returns path."""
    print(f"  Converting {ext.upper()[1:]} → Markdown …")
    if ext == ".docx":
        md_text = convert_docx_to_md(input_path)
    elif ext == ".pptx":
        md_text = convert_pptx_to_md(input_path)
    elif ext == ".xlsx":
        md_text = convert_xlsx_to_md(input_path)
    else:
        raise ValueError(f"Cannot convert {ext}")

    tmp = tempfile.NamedTemporaryFile(
        mode="w", suffix=".md", delete=False, encoding="utf-8",
        prefix=f"pageindex_{Path(input_path).stem}_"
    )
    tmp.write(md_text)
    tmp.close()
    return tmp.name


def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="pageindex-cli",
        description="Convert documents to PageIndex JSON node tree using Qwen (or any OpenAI-compat) API.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # PDF with Qwen API
  pageindex-cli doc.pdf --qwen-api-key sk-xxx --qwen-model qwen-long

  # DOCX with custom endpoint
  pageindex-cli report.docx \\
      --qwen-api-key sk-xxx \\
      --qwen-api-base https://dashscope.aliyuncs.com/compatible-mode/v1 \\
      --qwen-model qwen-plus

  # PPTX, save to specific file
  pageindex-cli slides.pptx --qwen-api-key sk-xxx -o slides_index.json

  # XLSX without summaries
  pageindex-cli data.xlsx --qwen-api-key sk-xxx --no-summary --no-description

  # Use env vars instead of CLI flags
  export QWEN_API_KEY=sk-xxx
  export QWEN_API_BASE=https://dashscope.aliyuncs.com/compatible-mode/v1
  export QWEN_MODEL=qwen-long
  pageindex-cli document.pdf
""",
    )

    # ── Positional ──
    p.add_argument("input", metavar="FILE", help="Input file (PDF/MD/DOCX/PPTX/XLSX)")

    # ── Output ──
    p.add_argument("-o", "--output", metavar="FILE",
                   help="Output JSON path (default: <input_stem>_index.json)")

    # ── Qwen / API config ──
    api = p.add_argument_group("API Configuration")
    api.add_argument("--qwen-api-key", metavar="KEY",
                     default=os.environ.get("QWEN_API_KEY") or os.environ.get("OPENAI_API_KEY"),
                     help="Qwen API key (or set QWEN_API_KEY env var)")
    api.add_argument("--qwen-api-base", metavar="URL",
                     default=os.environ.get("QWEN_API_BASE",
                                            "https://dashscope.aliyuncs.com/compatible-mode/v1"),
                     help="Qwen API base URL (default: DashScope compatible endpoint)")
    api.add_argument("--qwen-model", metavar="MODEL",
                     default=os.environ.get("QWEN_MODEL", "qwen-long"),
                     help="Model name (default: qwen-long)")

    # ── PageIndex options ──
    pi = p.add_argument_group("PageIndex Options")
    pi.add_argument("--toc-check-pages", type=int, default=20, metavar="N",
                    help="Pages to scan for TOC (PDF only, default: 20)")
    pi.add_argument("--max-pages-per-node", type=int, default=10, metavar="N",
                    help="Max pages per tree node (PDF only, default: 10)")
    pi.add_argument("--max-tokens-per-node", type=int, default=20000, metavar="N",
                    help="Max tokens per tree node (PDF only, default: 20000)")
    pi.add_argument("--no-summary", action="store_true",
                    help="Disable per-node summaries (faster, smaller output)")
    pi.add_argument("--no-description", action="store_true",
                    help="Disable document-level description")
    pi.add_argument("--include-text", action="store_true",
                    help="Include raw text in each node (larger output)")
    pi.add_argument("--no-node-id", action="store_true",
                    help="Omit node IDs from output")

    # ── MD-specific ──
    md = p.add_argument_group("Markdown / Office Options")
    md.add_argument("--thinning", action="store_true",
                    help="Apply tree thinning for Markdown (merge small nodes)")
    md.add_argument("--thinning-threshold", type=int, default=5000, metavar="N",
                    help="Min token count before thinning (default: 5000)")
    md.add_argument("--summary-token-threshold", type=int, default=200, metavar="N",
                    help="Min tokens in node before generating summary (default: 200)")

    # ── Misc ──
    p.add_argument("--keep-temp", action="store_true",
                   help="Keep temporary MD file created from DOCX/PPTX/XLSX")
    p.add_argument("--pretty", action="store_true", default=True,
                   help="Pretty-print JSON output (default: true)")
    p.add_argument("--compact", dest="pretty", action="store_false",
                   help="Compact JSON output (no indentation)")

    return p


def main():
    parser = build_parser()
    args = parser.parse_args()

    print_banner()

    # ── Validate input ──
    input_path = os.path.abspath(args.input)
    if not os.path.isfile(input_path):
        sys.exit(f"File not found: {input_path}")

    ext = detect_format(input_path)
    print(f"  Input  : {input_path}")
    print(f"  Format : {ext.upper()[1:]}")

    # ── API key check ──
    if not args.qwen_api_key:
        sys.exit(
            "No API key provided.\n"
            "Pass --qwen-api-key or set the QWEN_API_KEY environment variable."
        )

    # ── Configure environment for LiteLLM ──
    setup_qwen_env(args.qwen_api_key, args.qwen_api_base)
    print(f"  Model  : {args.qwen_model}")
    print(f"  API URL: {args.qwen_api_base}")

    # ── Output path ──
    stem = Path(input_path).stem
    output_path = args.output or os.path.join(
        os.path.dirname(input_path), f"{stem}_index.json"
    )
    output_path = os.path.abspath(output_path)
    print(f"  Output : {output_path}")
    print()

    # ── Convert office formats → temp MD ──
    temp_md_path = None
    process_as_pdf = (ext == ".pdf")
    process_as_md_path = input_path if ext in (".md", ".markdown") else None

    if ext in (".docx", ".pptx", ".xlsx"):
        temp_md_path = convert_to_temp_md(input_path, ext)
        process_as_md_path = temp_md_path
        print(f"  Temp MD: {temp_md_path}")
        print()

    # ── Build PageIndex opts ──
    opts = {
        "toc_check_pages": args.toc_check_pages,
        "max_pages_per_node": args.max_pages_per_node,
        "max_tokens_per_node": args.max_tokens_per_node,
        "if_add_node_id": "no" if args.no_node_id else "yes",
        "if_add_node_summary": "no" if args.no_summary else "yes",
        "if_add_doc_description": "no" if args.no_description else "yes",
        "if_add_node_text": "yes" if args.include_text else "no",
        "if_thinning": args.thinning,
        "thinning_threshold": args.thinning_threshold,
        "summary_token_threshold": args.summary_token_threshold,
    }

    # ── Run PageIndex ──
    # LiteLLM uses "openai/<model>" for custom base URLs
    litellm_model = f"openai/{args.qwen_model}"

    try:
        if process_as_pdf:
            print("  Processing PDF with PageIndex …")
            result = run_pageindex_pdf(input_path, litellm_model, opts)
        else:
            print("  Processing Markdown with PageIndex …")
            result = run_pageindex_md(process_as_md_path, litellm_model, opts)
    except Exception as e:
        sys.exit(f"\nPageIndex error: {e}")
    finally:
        if temp_md_path and not args.keep_temp:
            try:
                os.unlink(temp_md_path)
            except OSError:
                pass

    # ── Attach source metadata ──
    if isinstance(result, dict):
        result["_meta"] = {
            "source_file": input_path,
            "source_format": ext.lstrip(".").upper(),
            "model": args.qwen_model,
            "api_base": args.qwen_api_base,
        }

    # ── Save JSON ──
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    indent = 2 if args.pretty else None
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, indent=indent, ensure_ascii=False)

    print()
    print(f"✅  Done! Tree structure saved to:\n    {output_path}")

    # ── Quick stats ──
    def count_nodes(obj):
        if not isinstance(obj, (dict, list)):
            return 0
        if isinstance(obj, list):
            return sum(count_nodes(i) for i in obj)
        total = 1
        for v in obj.values():
            total += count_nodes(v)
        return total

    structure = result.get("structure", [])
    node_count = count_nodes(structure)
    print(f"    Nodes  : {node_count}")
    print(f"    Size   : {os.path.getsize(output_path):,} bytes")
    print()


if __name__ == "__main__":
    main()
